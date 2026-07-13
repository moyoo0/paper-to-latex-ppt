#!/usr/bin/env python3
"""Convert a final slide PDF into a PPTX with one full-slide image per page."""

from __future__ import annotations

import argparse
import tempfile
from pathlib import Path


def require_dependencies():
    try:
        import fitz  # type: ignore
    except ImportError as exc:
        raise SystemExit(
            "Missing dependency: PyMuPDF. Install with `python3 -m pip install pymupdf python-pptx`."
        ) from exc
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches  # type: ignore
    except ImportError as exc:
        raise SystemExit(
            "Missing dependency: python-pptx. Install with `python3 -m pip install python-pptx`."
        ) from exc
    return fitz, Presentation, Inches


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pdf", type=Path)
    parser.add_argument("--out", type=Path, default=Path("output/final.pptx"))
    parser.add_argument("--template-pptx", type=Path)
    parser.add_argument("--dpi", type=int, default=300)
    args = parser.parse_args()

    fitz, Presentation, Inches = require_dependencies()
    doc = fitz.open(args.pdf)
    if doc.page_count == 0:
        raise SystemExit("PDF has no pages.")

    prs = Presentation(str(args.template_pptx)) if args.template_pptx else Presentation()
    first = doc[0].rect
    width_in = first.width / 72
    height_in = first.height / 72
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    blank_layout = prs.slide_layouts[6]

    zoom = args.dpi / 72
    matrix = fitz.Matrix(zoom, zoom)
    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        for index, page in enumerate(doc, start=1):
            image_path = tmpdir / f"slide_{index:03d}.png"
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            pix.save(image_path)
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"Wrote {args.out}")


if __name__ == "__main__":
    main()
